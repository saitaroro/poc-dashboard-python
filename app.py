import os
import pandas as pd
import matplotlib
matplotlib.use('Agg') # compatilibite Docker pour les graphiques
import matplotlib.pyplot as pltate, re
from flask import Flask, render_template, send_file, redirect, url_for
from pptx import Presentation
from pptx.util import Inches
from email.message import EmailMessage
import io

# tentative d'import pour la génération de template Outlook (.oft)
try:
    import win32com.client
except ImportError:
    win32com = None

app = Flask(__name__)

# Configuration des dossiers
DATA_FILE = 'data/data_source.csv'
OUTPUT_DIR = 'output'
os.makedirs(OUTPUT_DIR, exist_ok=True)

# --- 1. Génération de fausses données pour le POC ---
def create_mock_data():
    # génère un jeu de données synthétique selon la nouvelle structure décrite par l'utilisateur
    if not os.path.exists(DATA_FILE):
        os.makedirs('data', exist_ok=True)
        rows = []
        start_date = datetime(2024, 1, 1)
        for i in range(500):
            # date du rendez-vous entre début 2024 et fin 2024
            dr = start_date + timedelta(days=random.randint(0, 365))
            # création du rdv jusqu'à 30 jours avant la date du rdv
            dc = dr - timedelta(days=random.randint(0, 30))
            canal = random.choice(['telephone', 'web'])
            profil = random.choice(['Junior', 'Senior', 'Manager', 'Client'])
            motif = random.choice(['Comptabilité', 'Support', 'Vente'])
            sous_motif = random.choice(['Question', 'Réclamation', 'Demande'])
            conseiller = random.randint(1, 20)
            bureau = random.choice(['Paris', 'Lyon', 'Marseille', 'Bordeaux'])
            mois = dr.strftime('%B')
            annee_mois = dr.strftime('%Y-%m')
            nb = random.randint(1, 3)  # agrégat
            rows.append({
                'date_rdv': dr.strftime('%Y-%m-%d'),
                'canal': canal,
                'profil': profil,
                'motif': motif,
                'sous_motif': sous_motif,
                'date_creation': dc.strftime('%Y-%m-%d'),
                'id_conseiller': conseiller,
                'bureau': bureau,
                'mois': mois,
                'annee_mois': annee_mois,
                'nb_rdv': nb
            })
        df = pd.DataFrame(rows)
        df.to_csv(DATA_FILE, index=False)

# --- 2. Fonctions de Calcul et Graphiques ---
def process_data():
    # lecture et parsing des dates
    df = pd.read_csv(DATA_FILE, parse_dates=['date_rdv', 'date_creation'])

    # volumes
    total_volume = df['nb_rdv'].sum()
    monthly = df.groupby('annee_mois')['nb_rdv'].sum()
    df['week'] = df['date_rdv'].dt.isocalendar().week
    weekly = df.groupby('week')['nb_rdv'].sum()

    # délai moyen en jours
    df['delay'] = (df['date_rdv'] - df['date_creation']).dt.days
    avg_delay = df['delay'].mean()

    # répartition journalière cumulée
    daily = df.groupby('date_rdv')['nb_rdv'].sum().cumsum()

    # fenêtre M-2 à M+1 par rapport au dernier mois disponible
    last = df['date_rdv'].max().replace(day=1)
    m2 = last - pd.DateOffset(months=2)
    m_plus1 = last + pd.DateOffset(months=1)
    mask = (df['date_rdv'] >= m2) & (df['date_rdv'] <= m_plus1)
    window_vol = df.loc[mask].groupby('annee_mois')['nb_rdv'].sum()

    # motifs et sous-motifs
    motif_counts = df.groupby('motif')['nb_rdv'].sum()
    sous_motif_counts = df.groupby('sous_motif')['nb_rdv'].sum()

    # régions dérivées du bureau (simplification)
    region_map = {
        'Paris': "Île-de-France",
        'Lyon': "Auvergne-Rhône-Alpes",
        'Marseille': "Provence-Alpes-Côte d'Azur",
        'Bordeaux': "Nouvelle-Aquitaine"
    }
    df['region'] = df['bureau'].map(region_map)
    region_vol = df.groupby('region')['nb_rdv'].sum()

    # évolution mois à mois par région
    df['annee_mois_dt'] = pd.to_datetime(df['annee_mois'] + '-01')
    region_monthly = df.groupby(['region', 'annee_mois'])['nb_rdv'].sum().unstack(fill_value=0)
    months = sorted(region_monthly.columns)
    if len(months) >= 2:
        prev, lastm = months[-2], months[-1]
        region_mom = region_monthly[lastm] - region_monthly[prev]
    else:
        region_mom = pd.Series(dtype=float)

    # tendance globale (% change)
    trend = monthly.pct_change().fillna(0)

    # générer quelques graphiques de base (pour démonstration)
    charts = {}
    plt.figure()
    monthly.plot(kind='bar', color='skyblue')
    plt.title('Volumes mensuels')
    plt.xlabel('Année-Mois')
    plt.ylabel('Rendez-vous')
    plt.tight_layout()
    charts['monthly'] = os.path.join(OUTPUT_DIR, 'monthly.png')
    plt.savefig(charts['monthly'])
    plt.close()

    plt.figure()
    df['delay'].hist(bins=20)
    plt.title('Délai de traitement (jours)')
    plt.tight_layout()
    charts['delay'] = os.path.join(OUTPUT_DIR, 'delay.png')
    plt.savefig(charts['delay'])
    plt.close()

    plt.figure(figsize=(10, 6))
    region_monthly.T.plot()
    plt.title('Évolution par région (mois)')
    plt.tight_layout()
    charts['region_trend'] = os.path.join(OUTPUT_DIR, 'region_trend.png')
    plt.savefig(charts['region_trend'])
    plt.close()

    return {
        'total_volume': total_volume,
        'monthly': monthly,
        'weekly': weekly,
        'avg_delay': avg_delay,
        'daily_cumulative': daily,
        'window_vol': window_vol,
        'motif_counts': motif_counts,
        'sous_motif_counts': sous_motif_counts,
        'region_vol': region_vol,
        'region_mom': region_mom,
        'trend': trend,
        'charts': charts
    }

# --- 3. Génération du PowerPoint ---
def generate_pptx(results: dict):
    """Crée un fichier PowerPoint à partir des indicateurs calculés.

    * `results` est le dictionnaire renvoyé par `process_data`.
    * On inclut les graphiques pré‑générés et quelques métriques clés.
    """
    prs = Presentation()

    # Slide 1 : titre général
    slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    subtitle = slide.placeholders[1]
    title.text = "Rapport Mensuel 2025"
    subtitle.text = "Analyse des rendez‑vous"

    # Slide 2 : aperçu des volumes
    slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(slide_layout)
    slide.shapes.title.text = "Vue d'ensemble"
    tf = slide.shapes.placeholders[1].text_frame
    tf.text = f"Volume total de rendez‑vous : {results['total_volume']}"
    p = tf.add_paragraph()
    p.text = f"Délai moyen de traitement : {results['avg_delay']:.1f} jours"

    # Insérer les graphiques générés automatiquement
    y_offset = 3.0
    for name, path in results['charts'].items():
        try:
            slide.shapes.add_picture(path, Inches(1), Inches(y_offset), width=Inches(8))
            y_offset += 4.5
            # ajouter une slide si l'espace manque
            if y_offset > 7:
                slide_layout = prs.slide_layouts[1]
                slide = prs.slides.add_slide(slide_layout)
                y_offset = 1.0
        except Exception:
            # si un graphique ne peut pas être ajouté, on l'ignore
            pass

    pptx_path = os.path.join(OUTPUT_DIR, 'Rapport_Final.pptx')
    prs.save(pptx_path)
    return pptx_path

# --- ROUTES FLASK (Le Site Web) ---

@app.route('/')
def index():
    create_mock_data() # S'assure que le fichier existe
    return render_template('index.html')

@app.route('/integrate', methods=['POST'])
def integrate():
    # Lance les calculs et récupère tous les indicateurs
    stats = process_data()
    # Génère le PPT à partir du dictionnaire renvoyé
    generate_pptx(stats)
    return redirect(url_for('validation_page'))

@app.route('/validate')
def validation_page():
    return render_template('validate.html')

@app.route('/test-download')
def test_download():
    # Permet de télécharger le PPT pour vérifier
    return send_file(os.path.join(OUTPUT_DIR, 'Rapport_Final.pptx'), as_attachment=True)

@app.route('/generate-email', methods=['POST'])
def generate_email():
    # Récupération des données du formulaire
    email_to = request.form.get('email_to')
    email_cc = request.form.get('email_cc')
    subject = request.form.get('subject')

    body_text = (
        "Bonjour,\n\nVeuillez trouver ci-joint le rapport mensuel des rendez-vous.\n\nCordialement."
    )

    pptx_path = os.path.join(OUTPUT_DIR, 'Rapport_Final.pptx')

    # si win32com est disponible, on génère un template .oft via Outlook
    if win32com is not None:
        try:
            outlook = win32com.client.Dispatch('Outlook.Application')
            mail = outlook.CreateItem(0)  # 0 == olMailItem
            mail.Subject = subject
            mail.To = email_to
            mail.CC = email_cc
            mail.Body = body_text
            mail.Attachments.Add(pptx_path)

            oft_path = os.path.join(OUTPUT_DIR, 'template.oft')
            # 5 == olTemplate
            mail.SaveAs(oft_path, 5)
            return send_file(oft_path, as_attachment=True, download_name="template.oft")
        except Exception as e:
            # si Outlook n'est pas installé ou erreur, on retombe sur .eml
            print(f"Erreur génération OFT : {e}")

    # fallback classique en .eml
    msg = EmailMessage()
    msg['Subject'] = subject
    msg['From'] = "monapp@entreprise.com"
    msg['To'] = email_to
    msg['Cc'] = email_cc
    msg.set_content(body_text)
    with open(pptx_path, 'rb') as f:
        file_data = f.read()
        msg.add_attachment(
            file_data,
            maintype='application',
            subtype='vnd.openxmlformats-officedocument.presentationml.presentation',
            filename="Rapport_2025.pptx",
        )

    eml_path = os.path.join(OUTPUT_DIR, 'brouillon_outlook.eml')
    with open(eml_path, 'wb') as f:
        f.write(msg.as_bytes())

    return send_file(eml_path, as_attachment=True, download_name="Ouvrir_dans_Outlook.eml")

if __name__ == '__main__':
    app.run(host='0.0.0.0', port=5000, debug=True)